from langchain.schema import HumanMessage, SystemMessage
import os
from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
import pandas as pd
import altair as alt
import plotly.express as px
import streamlit as st
import tempfile
import base64
import io
import numpy as np
from langchain_experimental.agents.agent_toolkits import create_pandas_dataframe_agent
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.chains import RetrievalQA
import sqlite3
from sqlalchemy import create_engine, text
import openpyxl
from youtube_transcript_api import YouTubeTranscriptApi
import re
from PIL import Image
import google.generativeai as genai
import docx2txt
from pptx import Presentation
import json
import requests
from urllib.parse import urlparse, parse_qs
import supabase
from supabase import create_client, Client
import logging
from datetime import datetime, timedelta
import time
from functools import wraps
import hashlib
from collections import defaultdict
from typing import Optional, Dict, Any
import shutil

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('app.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Constants
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
ALLOWED_EXTENSIONS = {
    'csv': ['csv'],
    'documents': ['pdf', 'docx', 'pptx'],
    'images': ['jpg', 'jpeg', 'png'],
    'spreadsheets': ['xlsx', 'xls'],
    'databases': ['sql', 'db']
}
RATE_LIMIT_CALLS = 10
RATE_LIMIT_PERIOD = 60
SESSION_TIMEOUT = 30

supabase_url = os.getenv("SUPABASE_URL")
supabase_key = os.getenv("SUPABASE_KEY")
if not supabase_url or not supabase_key:
    raise ValueError("Missing Supabase credentials")
supabase: Client = create_client(supabase_url=supabase_url, supabase_key=supabase_key)

class RateLimiter:
    def __init__(self, max_calls: int, period: int):
        self.max_calls = max_calls
        self.period = period
        self.calls = defaultdict(list)
    
    def is_allowed(self, key: str) -> bool:
        now = time.time()
        self.calls[key] = [call_time for call_time in self.calls[key] if now - call_time < self.period]
        
        if len(self.calls[key]) >= self.max_calls:
            return False
        
        self.calls[key].append(now)
        return True

rate_limiter = RateLimiter(max_calls=10, period=60)  # 10 calls per minute

def check_rate_limit(key: str) -> bool:
    """Check if the request is within rate limits."""
    return rate_limiter.is_allowed(key)

def validate_file(file) -> bool:
    """Validate uploaded file size and type."""
    if file is None:
        return False
    
    # Check file size
    file_size = 0
    for chunk in file:
        file_size += len(chunk)
    if file_size > MAX_FILE_SIZE:
        return False
    
    file_ext = file.name.split('.')[-1].lower()
    allowed_extensions = [ext for exts in ALLOWED_EXTENSIONS.values() for ext in exts]
    return file_ext in allowed_extensions

def generate_csrf_token() -> str:
    """Generate a CSRF token."""
    if 'csrf_token' not in st.session_state:
        st.session_state.csrf_token = hashlib.sha256(str(time.time()).encode()).hexdigest()
    return st.session_state.csrf_token

def verify_csrf_token(token: str) -> bool:
    """Verify CSRF token."""
    return token == st.session_state.get('csrf_token')

def check_session_timeout():
    """Check if the current session has timed out."""
    if 'last_activity' in st.session_state:
        last_activity = datetime.fromisoformat(st.session_state.last_activity)
        if datetime.now() - last_activity > timedelta(minutes=SESSION_TIMEOUT):
            sign_out()
            st.error("Your session has timed out. Please log in again.")
            st.rerun()
    st.session_state.last_activity = datetime.now().isoformat()

class AppError(Exception):
    """Base exception for application errors."""
    pass

class AuthenticationError(AppError):
    """Raised when authentication fails."""
    pass

class PaymentError(AppError):
    """Raised when payment processing fails."""
    pass

class FileProcessingError(AppError):
    """Raised when file processing fails."""
    pass

def handle_error(func):
    """Decorator for handling errors in functions."""
    @wraps(func)
    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except AppError as e:
            logger.error(f"Application error in {func.__name__}: {str(e)}")
            st.error(str(e))
        except Exception as e:
            logger.error(f"Unexpected error in {func.__name__}: {str(e)}")
            st.error("An unexpected error occurred. Please try again later.")
    return wrapper

def save_payment_status(user_id: str, session_id: str, status: str):
    """Save payment status to database."""
    try:
        data = {
            'user_id': user_id,
            'session_id': session_id,
            'status': status,
            'created_at': datetime.now().isoformat()
        }
        supabase.table('payments').insert(data).execute()
    except Exception as e:
        logger.error(f"Error saving payment status: {str(e)}")
        raise PaymentError("Failed to save payment status")

def get_payment_status(user_id: str) -> Optional[Dict[str, Any]]:
    """Get payment status from database."""
    try:
        response = supabase.table('payments')\
            .select('*')\
            .eq('user_id', user_id)\
            .order('created_at', desc=True)\
            .limit(1)\
            .execute()
        return response.data[0] if response.data else None
    except Exception as e:
        logger.error(f"Error getting payment status: {str(e)}")
        return None

@handle_error
def sign_up(email: str, password: str):
    """Handle user registration with proper validation and error handling."""
    try:
        if not email or not password:
            raise AuthenticationError("Please provide both email and password")
        
        if len(password) < 8:
            raise AuthenticationError("Password must be at least 8 characters long")
        
        # Rate limiting check
        if not check_rate_limit(f"signup_{email}"):
            raise AuthenticationError("Too many signup attempts. Please try again later.")
        
        user = supabase.auth.sign_up({"email": email, "password": password})
        if user:
            logger.info(f"New user registered: {email}")
            st.success("Registration successful! Please check your email for a confirmation link.")
            st.info("After confirming your email, you can log in to your account.")
        return user
    
    except Exception as e:
        error_message = str(e)
        if "security purposes" in error_message:
            logger.warning(f"Rate limit hit for registration: {email}")
            st.warning("Please wait a moment before trying again. This is a security measure to protect your account.")
        else:
            logger.error(f"Registration failed for {email}: {str(e)}")
        st.error(f"Registration Failed: {e}")
        return None

@handle_error
def sign_in(email: str, password: str):
    """Handle user login with proper validation and error handling."""
    try:
        if not email or not password:
            raise AuthenticationError("Please provide both email and password")

        if not check_rate_limit(f"login_{email}"):
            raise AuthenticationError("Too many login attempts. Please try again later.")
        
        user = supabase.auth.sign_in_with_password({"email": email, "password": password})
        if user:
            st.session_state.user_email = email
            st.session_state.user_id = user.user.id
            st.session_state.last_activity = datetime.now().isoformat()
            logger.info(f"User logged in: {email}")
            st.success(f"Welcome back, {email}")
            st.rerun()
        return user
    
    except Exception as e:
        logger.error(f"Login failed for {email}: {str(e)}")
        st.error(f"Login Failed: {e}")
        return None

def sign_out():
    try:
        supabase.auth.sign_out()
        st.session_state.user_email = None
        st.rerun()

    except Exception as e:
        st.error(f"Logout Failed: {e}")

def main_app():
    """Main application with proper security and error handling."""
    check_session_timeout()

    try:
        llm = ChatGoogleGenerativeAI(
            model="gemini-1.5-flash",
            temperature=0,
            max_tokens=None,
            timeout=None,
            max_retries=2,
        )
    except Exception as e:
        logger.error(f"Error initializing LLM: {str(e)}")
        st.error("Error initializing AI model. Please try again later.")
        return

    CSV_PROMPT_PREFIX = "You are a data analysis expert. Please analyze the following data and answer the question: "
    CSV_PROMPT_SUFFIX = "\nProvide a clear and concise answer based on the data."

    CODE_PROMPT_PREFIX = """
    You are an expert code writer and reviewer. Your task is to:
    1. Write clean, efficient, error-less, and well-documented code
    2. Review and improve existing code
    3. Explain code functionality
    4. Suggest optimizations and best practices
    5. Fix bugs and issues
    6. Add error handling and input validation
    7. Implement security best practices

    Please analyze the following code or request:
    """

    CODE_PROMPT_SUFFIX = """
    Provide your response in the following format:
    1. Code Analysis (if reviewing existing code)
    2. Implementation/Improvements
    3. Explanation
    4. Best Practices Applied
    5. Security Considerations
    6. Error Handling
    7. Testing Recommendations

    Use markdown formatting for better readability.
    """

    # Initialize session state variables at the start
    if "user_email" not in st.session_state:
        st.session_state.user_email = None
    if "question" not in st.session_state:
        st.session_state.question = ""

    def analyze_and_visualize(df):
        st.write("### Data Analysis and Visualizations")

        st.write("#### Basic Statistics")
        st.write(df.describe())
        
        tab1, tab2, tab3 = st.tabs(["Numerical Analysis", "Categorical Analysis", "Correlation Analysis"])
        
        with tab1:
            st.write("#### Numerical Columns Distribution")
            numerical_cols = df.select_dtypes(include=[np.number]).columns
            if len(numerical_cols) > 0:
                selected_num_col = st.selectbox("Select numerical column", numerical_cols, key="num_col")
                fig = px.histogram(df, x=selected_num_col, title=f"Distribution of {selected_num_col}")
                st.plotly_chart(fig, use_container_width=True)

                fig_box = px.box(df, y=selected_num_col, title=f"Box Plot of {selected_num_col}")
                st.plotly_chart(fig_box, use_container_width=True)
        
        with tab2:
            st.write("#### Categorical Columns Analysis")
            categorical_cols = df.select_dtypes(include=['object']).columns
            if len(categorical_cols) > 0:
                selected_cat_col = st.selectbox("Select categorical column", categorical_cols, key="cat_col")
                value_counts = df[selected_cat_col].value_counts()
                fig = px.bar(x=value_counts.index, y=value_counts.values, 
                            title=f"Count of {selected_cat_col}",
                            labels={'x': selected_cat_col, 'y': 'Count'})
                st.plotly_chart(fig, use_container_width=True)

                fig_pie = px.pie(values=value_counts.values, names=value_counts.index, 
                            title=f"Distribution of {selected_cat_col}")
                st.plotly_chart(fig_pie, use_container_width=True)
        
        with tab3:
            st.write("#### Correlation Analysis")
            if len(numerical_cols) > 1:
                corr_matrix = df[numerical_cols].corr()
                fig = px.imshow(corr_matrix, 
                            title="Correlation Heatmap",
                            color_continuous_scale='RdBu_r')
                st.plotly_chart(fig, use_container_width=True)
                
                if len(numerical_cols) > 1:
                    st.write("#### Scatter Plot Matrix")
                    top_corr_pairs = []
                    for i in range(len(numerical_cols)):
                        for j in range(i+1, len(numerical_cols)):
                            corr = corr_matrix.iloc[i,j]
                            if abs(corr) > 0.5: 
                                top_corr_pairs.append((numerical_cols[i], numerical_cols[j], corr))
                    
                    if top_corr_pairs:
                        selected_pair = st.selectbox("Select feature pair", 
                                                [f"{pair[0]} vs {pair[1]} (corr: {pair[2]:.2f})" 
                                                    for pair in top_corr_pairs])
                        selected_idx = [f"{pair[0]} vs {pair[1]} (corr: {pair[2]:.2f})" 
                                    for pair in top_corr_pairs].index(selected_pair)
                        x_col, y_col = top_corr_pairs[selected_idx][0], top_corr_pairs[selected_idx][1]
                        
                        fig_scatter = px.scatter(df, x=x_col, y=y_col, 
                                            title=f"Scatter Plot: {x_col} vs {y_col}")
                        st.plotly_chart(fig_scatter, use_container_width=True)

    def process_csv(file):
        df = pd.read_csv(file, encoding="utf-8")
        agent = create_pandas_dataframe_agent(
            llm=llm,
            df=df,
            verbose=True,
            allow_dangerous_code=True,
            tools=[analyze_and_visualize]
        )
        return df, agent

    def process_pdf(file):
        with open("temp.pdf", "wb") as f:
            f.write(file.getvalue())

        loader = PyPDFLoader("temp.pdf")
        pages = loader.load()

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        chunks = text_splitter.split_documents(pages)

        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        vectorstore = FAISS.from_documents(chunks, embeddings)

        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=vectorstore.as_retriever()
        )
        return qa_chain

    def process_parquet(file):
        df = pd.read_parquet(file)
        agent = create_pandas_dataframe_agent(
            llm=llm,
            df=df,
            verbose=True,
            allow_dangerous_code=True,
            tools=[analyze_and_visualize]
        )
        return df, agent

    def process_excel(file):
        excel_file = pd.ExcelFile(file)
        sheet_names = excel_file.sheet_names
        
        dfs = {}
        for sheet in sheet_names:
            dfs[sheet] = pd.read_excel(file, sheet_name=sheet)
    
        agents = {}
        for sheet, df in dfs.items():
            agents[sheet] = create_pandas_dataframe_agent(
                llm=llm,
                df=df,
                verbose=True,
                allow_dangerous_code=True,
                tools=[analyze_and_visualize]
            )
        
        return dfs, agents

    def process_sql(file):
        try:
            file_content = file.read()
            file.seek(0)
            
            if file_content.startswith(b'SQLite format 3'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.db') as temp_db:
                    temp_db.write(file_content)
                    temp_db_path = temp_db.name
                
                conn = sqlite3.connect(temp_db_path)
            else:
                conn = sqlite3.connect(':memory:')
                file_content = file_content.replace(b'\x00', b'')
                try:
                    sql_script = file_content.decode('utf-8')
                except UnicodeDecodeError:
                    sql_script = file_content.decode('latin-1')
                
                statements = sql_script.split(';')
                
                for statement in statements:
                    statement = statement.strip()
                    if statement:
                        try:
                            conn.execute(statement)
                        except sqlite3.Error as e:
                            st.error(f"SQL Error in statement: {statement[:100]}...\nError: {str(e)}")
                            return None, None
            
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            
            if not tables:
                st.warning("No tables found in the database. Please check if the file is valid.")
                return None, None
            
            dfs = {}
            for table in tables:
                table_name = table[0]
                try:
                    dfs[table_name] = pd.read_sql_query(f"SELECT * FROM {table_name}", conn)
                except Exception as e:
                    st.warning(f"Could not read table {table_name}: {str(e)}")
                    continue
            
            if not dfs:
                st.error("No valid tables could be read from the database.")
                return None, None
            
            agents = {}
            for table, df in dfs.items():
                agents[table] = create_pandas_dataframe_agent(
                    llm=llm,
                    df=df,
                    verbose=True,
                    allow_dangerous_code=True,
                    tools=[analyze_and_visualize]
                )
            
            st.success("Database file successfully uploaded and processed!")
            
            selected_table = st.selectbox("Select Table", list(dfs.keys()))
            df = dfs[selected_table]
            agent = agents[selected_table]
            
            st.write("### Dataset Preview")
            st.write(df.head())
            
            if 'temp_db_path' in locals():
                try:
                    os.unlink(temp_db_path)
                except:
                    pass
            
            return df, agent
            
        except Exception as e:
            st.error(f"Error processing database file: {str(e)}")
            if 'temp_db_path' in locals():
                try:
                    os.unlink(temp_db_path)
                except:
                    pass
            return None, None

    def extract_video_id(url):
        patterns = [
            r'(?:youtube\.com\/watch\?v=|youtu\.be\/)([^&\n?]+)',
            r'youtube\.com\/embed\/([^&\n?]+)',
            r'youtube\.com\/v\/([^&\n?]+)'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        return None

    def get_transcript(video_id):
        try:
            # First try to get English transcript
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en'])
                return ' '.join([entry['text'] for entry in transcript])
            except:
                # If English transcript is not available, try to get auto-generated transcript
                transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                
                try:
                    transcript = transcript_list.find_generated_transcript(['en']).fetch()
                    return ' '.join([entry['text'] for entry in transcript])
                except:
                    try:
                        available_transcripts = transcript_list.find_manually_created_transcript(['en'])
                        st.warning("English transcript is not available. The video has manually created transcripts in other languages.")
                    except:
                        st.warning("English transcript is not available. The video only has auto-generated transcripts in other languages.")

                    st.write("Available languages:")
                    for transcript in transcript_list:
                        st.write(f"- {transcript.language} ({transcript.language_code})")
                    
                    return None
        except Exception as e:
            st.error(f"Error getting transcript: {str(e)}")
            return None

    def summarize_video(transcript, custom_prompt=None):
        if not transcript:
            st.error("Cannot summarize video: No transcript is available in English.")
            return None
        
        prompt = custom_prompt or """
        Please provide a comprehensive summary of this video transcript. Include:
        1. Main topics discussed
        2. Key points and insights
        3. Important conclusions or takeaways
        4. Any notable quotes or statistics
        
        Format the summary in a clear, structured way using markdown.
        """
        
        response = llm.invoke(prompt + "\n\nTranscript:\n" + transcript)
        return response.content

    def summarize_image(image):
        try:
            model = genai.GenerativeModel('gemini-1.5-flash')
            
            if not isinstance(image, Image.Image):
                image = Image.open(io.BytesIO(image))
            
            img_byte_arr = io.BytesIO()
            image.save(img_byte_arr, format=image.format or 'JPEG')
            img_byte_arr = img_byte_arr.getvalue()
            
            image_parts = [
                {
                    "mime_type": "image/jpeg",
                    "data": img_byte_arr
                }
            ]

            response = model.generate_content([
                "Please provide a detailed description of this image. Include:",
                "1. Main subject and setting",
                "2. Colors and visual elements",
                "3. Any text or symbols visible",
                "4. Overall mood or atmosphere",
                "5. Any notable details or patterns",
                image_parts[0]
            ])
            
            return response.text
        except Exception as e:
            st.error(f"Error analyzing image: {str(e)}")
            return None

    def analyze_code(code, task=None, language=None):
        try:
            prompt = CODE_PROMPT_PREFIX
            if task:
                prompt += f"\nTask: {task}\n"
            if language:
                prompt += f"\nLanguage: {language}\n"
            prompt += f"\nCode:\n{code}\n" + CODE_PROMPT_SUFFIX
            
            response = llm.invoke(prompt)
            return response.content
        except Exception as e:
            st.error(f"Error analyzing code: {str(e)}")
            return None

    def generate_code(task, language="python"):
        try:
            prompt = f"""
            Write code in {language} for the following task:
            {task}
            
            Requirements:
            1. Clean and efficient code following {language} best practices and conventions
            2. Proper error handling and input validation
            3. Security best practices specific to {language}
            4. Clear documentation and comments
            5. Follow {language} style guide and conventions
            6. Include usage examples and test cases
            7. Consider language-specific features and idioms
            
            Provide the code with:
            1. Implementation
            2. Usage examples
            3. Explanation of key concepts
            4. Dependencies and requirements
            5. Testing instructions
            6. Common pitfalls and solutions
            
            Format the response in markdown with proper code blocks and syntax highlighting.
            """
            
            response = llm.invoke(prompt)
            return response.content
        except Exception as e:
            st.error(f"Error generating code: {str(e)}")
            return None

    def optimize_code(code, optimization_type="performance", language=None):
        try:
            prompt = f"""
            Optimize the following {language if language else ''} code for {optimization_type}:
            
            {code}
            
            Focus on:
            1. Code efficiency and performance
            2. Memory usage and resource management
            3. Algorithm complexity and optimization
            4. Language-specific best practices
            5. Code maintainability and readability
            6. Security considerations
            7. Error handling and edge cases
            
            Provide:
            1. Optimized code
            2. Explanation of improvements
            3. Performance impact analysis
            4. Alternative approaches
            5. Testing recommendations
            
            Format the response in markdown with proper code blocks and syntax highlighting.
            """
            
            response = llm.invoke(prompt)
            return response.content
        except Exception as e:
            st.error(f"Error optimizing code: {str(e)}")
            return None

    def extract_text_from_pptx(file):
        try:
            prs = Presentation(file)
            text_content = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"\n--- Slide {slide_number} ---\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                text_content.append("\n".join(slide_text))
            
            return "\n".join(text_content)
        except Exception as e:
            st.error(f"Error extracting text from PowerPoint: {str(e)}")
            return None

    def analyze_presentation_content(text, analysis_type="general"):
        try:
            prompt = f"""
            Analyze the following presentation content with focus on {analysis_type}:
            
            {text}
            
            Provide a comprehensive analysis including:
            1. Presentation Structure
            - Number of slides
            - Main sections
            - Flow and organization
            
            2. Content Analysis
            - Key messages and themes
            - Main points per slide
            - Supporting details
            
            3. Visual Elements
            - Text content
            - Potential visual elements (based on text descriptions)
            - Layout suggestions
            
            4. Presentation Quality
            - Clarity of message
            - Effectiveness of structure
            - Recommendations for improvement
            
            5. Target Audience
            - Identified audience
            - Appropriateness of content
            - Engagement level
            
            Format the analysis in markdown with clear sections and bullet points.
            """
            
            response = llm.invoke(prompt)
            return response.content
        except Exception as e:
            st.error(f"Error analyzing presentation: {str(e)}")
            return None

    def process_document(file, file_type):
        try:
            if file_type == 'pdf':
                with open("temp.pdf", "wb") as f:
                    f.write(file.getvalue())
                loader = PyPDFLoader("temp.pdf")
            elif file_type == 'docx':
                with open("temp.docx", "wb") as f:
                    f.write(file.getvalue())
                loader = Docx2txtLoader("temp.docx")
            elif file_type == 'pptx':
                with open("temp.pptx", "wb") as f:
                    f.write(file.getvalue())
                text = extract_text_from_pptx("temp.pptx")
                if text:
                    return text
                return None
            else:
                st.error(f"Unsupported document type: {file_type}")
                return None

            if file_type != 'pptx':
                pages = loader.load()
                
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200
                )
                chunks = text_splitter.split_documents(pages)

                embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
                vectorstore = FAISS.from_documents(chunks, embeddings)

                qa_chain = RetrievalQA.from_chain_type(
                    llm=llm,
                    chain_type="stuff",
                    retriever=vectorstore.as_retriever()
                )
                
                return qa_chain
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            st.error(f"Error processing document: {str(e)}")
            return None
        finally:
            try:
                if file_type == 'pdf' and os.path.exists("temp.pdf"):
                    os.remove("temp.pdf")
                elif file_type == 'docx' and os.path.exists("temp.docx"):
                    os.remove("temp.docx")
                elif file_type == 'pptx' and os.path.exists("temp.pptx"):
                    os.remove("temp.pptx")
            except:
                pass

    def extract_text_from_docx(file):
        try:
            text = docx2txt.process(file)
            return text
        except Exception as e:
            st.error(f"Error extracting text from Word document: {str(e)}")
            return None

    def analyze_document_content(text, analysis_type="general"):
        try:
            prompt = f"""
            Analyze the following document content with focus on {analysis_type}:
            
            {text}
            
            Provide a comprehensive analysis including:
            1. Main topics and themes
            2. Key points and arguments
            3. Important details and facts
            4. Structure and organization
            5. Writing style and tone
            6. Recommendations or conclusions
            
            Format the analysis in markdown with clear sections and bullet points.
            """
            
            response = llm.invoke(prompt)
            return response.content
        except Exception as e:
            st.error(f"Error analyzing document: {str(e)}")
            return None

    st.title("Document Analysis AI Agent")

    st.sidebar.title("Features")

    feature = st.sidebar.radio(
        "Select Feature",
        ["Document Analysis", "YouTube Summarizer", "Image Summarizer", "Code Writer"]
    )
    
    if feature == "Code Writer":
        st.header("AI Code Writer")
        
        code_task = st.selectbox(
            "Select Task",
            ["Write New Code", "Analyze Existing Code", "Optimize Code"]
        )
        
        if code_task == "Write New Code":
            language = st.text_input("Enter programming language (e.g., Python, JavaScript, Java, C++, Go, Rust, PHP, Ruby, Swift, Kotlin, TypeScript, etc.)")
            task_description = st.text_area("Describe what you want the code to do:", height=100)
            
            if st.button("Generate Code"):
                if task_description and language:
                    with st.spinner(f"Generating {language} code..."):
                        result = generate_code(task_description, language)
                        if result:
                            st.markdown("### Generated Code")
                            st.markdown(result)
                else:
                    st.warning("Please provide both a programming language and task description.")
                    
        elif code_task == "Analyze Existing Code":
            language = st.text_input("Enter programming language of the code")
            code = st.text_area("Paste your code here:", height=200)
            analysis_type = st.selectbox(
                "Analysis Type",
                ["General Review", "Security Analysis", "Performance Analysis", "Best Practices", "Code Quality", "Architecture Review"]
            )
            
            if st.button("Analyze Code"):
                if code:
                    with st.spinner("Analyzing code..."):
                        result = analyze_code(code, analysis_type, language if language else None)
                        if result:
                            st.markdown("### Code Analysis")
                            st.markdown(result)
                else:
                    st.warning("Please provide code to analyze.")
                    
        elif code_task == "Optimize Code":
            language = st.text_input("Enter programming language of the code")
            code = st.text_area("Paste your code here:", height=200)
            optimization_type = st.selectbox(
                "Optimization Type",
                ["Performance", "Memory Usage", "Code Readability", "Security", "Algorithm Efficiency", "Resource Management"]
            )
            
            if st.button("Optimize Code"):
                if code:
                    with st.spinner("Optimizing code..."):
                        result = optimize_code(code, optimization_type, language if language else None)
                        if result:
                            st.markdown("### Optimized Code")
                            st.markdown(result)
                else:
                    st.warning("Please provide code to optimize.")
    
    elif feature == "Document Analysis":
        uploaded_file = st.file_uploader(
            "Upload a file",
            type=["csv", "pdf", "parquet", "xlsx", "xls", "sql", "db", "docx", "pptx"],
            help="Upload your CSV, PDF, Parquet, Excel, SQL, Word, or PowerPoint file"
        )

        if uploaded_file is not None:
            # Validate file
            if not validate_file(uploaded_file):
                st.error(f"Invalid file. Please ensure the file is less than {MAX_FILE_SIZE/1024/1024}MB and has a supported format.")
                return

            try:
                file_type = uploaded_file.name.split('.')[-1].lower()
                logger.info(f"Processing file of type: {file_type}")
                
                # Create a temporary directory for file processing
                with tempfile.TemporaryDirectory() as temp_dir:
                    temp_file_path = os.path.join(temp_dir, uploaded_file.name)
                    logger.info(f"Created temporary file at: {temp_file_path}")
                    
                    # Save uploaded file to temporary directory
                    try:
                        with open(temp_file_path, 'wb') as f:
                            f.write(uploaded_file.getvalue())
                        logger.info("Successfully saved uploaded file")
                    except Exception as e:
                        logger.error(f"Error saving uploaded file: {str(e)}")
                        st.error("Error saving uploaded file. Please try again.")
                        return
                    
                    try:
                        if file_type in ['csv', 'parquet']:
                            if file_type == 'csv':
                                try:
                                    df = pd.read_csv(temp_file_path, encoding="utf-8")
                                    logger.info("Successfully read CSV file with utf-8 encoding")
                                except UnicodeDecodeError:
                                    df = pd.read_csv(temp_file_path, encoding="latin-1")
                                    logger.info("Successfully read CSV file with latin-1 encoding")
                                except Exception as e:
                                    logger.error(f"Error reading CSV file: {str(e)}")
                                    st.error("Error reading CSV file. Please ensure it's a valid CSV file.")
                                    return
                                
                                try:
                                    agent = create_pandas_dataframe_agent(
                                        llm=llm,
                                        df=df,
                                        verbose=True,
                                        allow_dangerous_code=True,
                                        tools=[analyze_and_visualize]
                                    )
                                    logger.info("Successfully created pandas dataframe agent")
                                except Exception as e:
                                    logger.error(f"Error creating pandas dataframe agent: {str(e)}")
                                    st.error("Error initializing data analysis. Please try again.")
                                    return
                            else:
                                try:
                                    df, agent = process_parquet(temp_file_path)
                                    logger.info("Successfully processed parquet file")
                                except Exception as e:
                                    logger.error(f"Error processing parquet file: {str(e)}")
                                    st.error("Error processing parquet file. Please ensure it's a valid parquet file.")
                                    return
                                
                            st.success(f"{file_type.upper()} file successfully uploaded and processed!")
                            st.write("### Dataset Preview")
                            st.write(df.head())
                            
                        elif file_type in ['pdf', 'docx', 'pptx']:
                            try:
                                result = process_document(uploaded_file, file_type)
                                logger.info(f"Successfully processed {file_type} document")
                            except Exception as e:
                                logger.error(f"Error processing {file_type} document: {str(e)}")
                                st.error(f"Error processing {file_type.upper()} file. Please ensure it's a valid {file_type.upper()} file.")
                                return
                                
                            if result:
                                st.success(f"{file_type.upper()} file successfully uploaded and processed!")
                                
                                if file_type in ['docx', 'pptx']:
                                    try:
                                        text = result if file_type == 'pptx' else extract_text_from_docx(uploaded_file)
                                        logger.info(f"Successfully extracted text from {file_type} file")
                                    except Exception as e:
                                        logger.error(f"Error extracting text from {file_type} file: {str(e)}")
                                        st.error(f"Error extracting text from {file_type.upper()} file. Please try again.")
                                        return
                                        
                                    if text:
                                        st.write("### Document Preview")
                                        st.text_area("Document Content", text[:1000] + "...", height=200)
                                        
                                        analysis_type = st.selectbox(
                                            "Analysis Type",
                                            ["General", "Technical", "Business", "Academic", "Legal", "Presentation"] if file_type == 'pptx' else ["General", "Technical", "Business", "Academic", "Legal"]
                                        )
                                        
                                        if st.button("Analyze Document"):
                                            with st.spinner("Analyzing document..."):
                                                try:
                                                    analysis = analyze_presentation_content(text, analysis_type) if file_type == 'pptx' else analyze_document_content(text, analysis_type)
                                                    logger.info(f"Successfully analyzed {file_type} document")
                                                except Exception as e:
                                                    logger.error(f"Error analyzing {file_type} document: {str(e)}")
                                                    st.error(f"Error analyzing {file_type.upper()} document. Please try again.")
                                                    return
                                                    
                                                if analysis:
                                                    st.markdown("### Document Analysis")
                                                    st.markdown(analysis)
                
                        elif file_type in ['xlsx', 'xls']:
                            try:
                                dfs, agents = process_excel(temp_file_path)
                                logger.info("Successfully processed Excel file")
                            except Exception as e:
                                logger.error(f"Error processing Excel file: {str(e)}")
                                st.error("Error processing Excel file. Please ensure it's a valid Excel file.")
                                return
                                
                            st.success("Excel file successfully uploaded and processed!")
                            
                            selected_sheet = st.selectbox("Select Sheet", list(dfs.keys()))
                            df = dfs[selected_sheet]
                            agent = agents[selected_sheet]
                            
                            st.write("### Dataset Preview")
                            st.write(df.head())
                            
                        elif file_type in ['sql', 'db']:
                            try:
                                df, agent = process_sql(temp_file_path)
                                logger.info("Successfully processed SQL/database file")
                            except Exception as e:
                                logger.error(f"Error processing SQL/database file: {str(e)}")
                                st.error("Error processing SQL/database file. Please ensure it's a valid database file.")
                                return
                                
                            if df is None or agent is None:
                                return
                        
                        st.write("### Ask a Question")
                        user_input = st.text_input(
                            "Enter your question:",
                            value=st.session_state.question,
                            key="question"
                        )

                        if st.button("Run Query"):
                            if user_input:
                                try:
                                    if file_type in ['csv', 'parquet', 'xlsx', 'xls', 'sql', 'db']:
                                        if any(keyword in user_input.lower() for keyword in ['visualize', 'visualization', 'plot', 'graph', 'chart', 'display visual', 'show me the data']):
                                            analyze_and_visualize(df)
                                        else:
                                            QUERY = CSV_PROMPT_PREFIX + user_input + CSV_PROMPT_SUFFIX
                                            res = agent.invoke(QUERY)
                                            st.write("### Final Answer")
                                            st.markdown(res["output"])
                                    elif file_type in ['pdf', 'docx']:
                                        if isinstance(result, RetrievalQA):
                                            res = result.invoke({"query": user_input})
                                            st.write("### Final Answer")
                                            st.markdown(res["result"])
                                        else:
                                            st.error("Document processing failed. Please try again.")
                                    elif file_type == 'pptx':
                                        if result:  # result is the text content
                                            analysis = analyze_presentation_content(result, "General")
                                            if analysis:
                                                st.markdown("### Document Analysis")
                                                st.markdown(analysis)
                                            else:
                                                st.error("Failed to analyze presentation. Please try again.")
                                except Exception as e:
                                    logger.error(f"Error processing query: {str(e)}")
                                    st.error("Error processing your question. Please try again with a different question.")
                    except Exception as e:
                        logger.error(f"Error in main file processing block: {str(e)}")
                        st.error("An error occurred while processing your file. Please try again.")
            except Exception as e:
                logger.error(f"Error in outer file processing block: {str(e)}")
                st.error("An error occurred while processing your file. Please try again.")
        else:
            st.info("Please upload a CSV, PDF, Parquet, Excel, SQL, Word, or PowerPoint file to get started")
            
    elif feature == "YouTube Summarizer":
        st.header("YouTube Video Summarizer")
        youtube_url = st.text_input("Enter YouTube URL")
        
        if youtube_url:
            video_id = extract_video_id(youtube_url)
            if video_id:
                st.video(youtube_url)
                
                if st.button("Get Summary"):
                    with st.spinner("Fetching transcript and generating summary..."):
                        transcript = get_transcript(video_id)
                        if transcript:
                            custom_prompt = st.text_area(
                                "Custom Summary Prompt (Optional)",
                                "Please provide a comprehensive summary of this video transcript. Include main topics, key points, and important takeaways.",
                                height=100
                            )
                            
                            summary = summarize_video(transcript, custom_prompt)
                            if summary:
                                st.markdown("### Video Summary")
                                st.markdown(summary)
            else:
                st.error("Invalid YouTube URL. Please enter a valid YouTube video URL.")
                
    elif feature == "Image Summarizer":
        st.header("Image Summarizer")
        uploaded_image = st.file_uploader(
            "Upload an image",
            type=["jpg", "jpeg", "png"],
            help="Upload an image to get a detailed description"
        )
        
        if uploaded_image:
            image = Image.open(uploaded_image)
            st.image(image, caption="Uploaded Image", use_column_width=True)
            
            if st.button("Analyze Image"):
                with st.spinner("Analyzing image..."):
                    description = summarize_image(image)
                    if description:
                        st.markdown("### Image Analysis")
                        st.markdown(description)

    if st.button("Logout"):
            sign_out()
            
def auth_screen():
    st.title("Streamlit + Supabase Authentication")
    option = st.selectbox("Choose an action:", ["Login", "Sign Up"])
    email = st.text_input("Email")
    password = st.text_input("Password", type="password")

    if option == "Sign Up" and st.button("Register"):
        user = sign_up(email, password)
        if user:
            st.success("Registration was successful. Please log in.")

    if option == "Login" and st.button("Login"):
        if not email or not password:
            st.error("Please enter both email and password")
            return
            
        user = sign_in(email, password)
        if user:
            st.session_state.user_email = email
            st.success(f"Welcome back, {email}")
            st.rerun()

if "user_email" not in st.session_state:
    st.session_state.user_email = None
if "question" not in st.session_state:
    st.session_state.question = ""

if st.session_state.user_email is not None:
    main_app()
else:
    auth_screen()
